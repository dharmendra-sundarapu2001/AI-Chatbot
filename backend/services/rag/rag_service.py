import os
import logging
import hashlib
from typing import List, Dict, Any, Optional
from pathlib import Path
import tempfile
import uuid
import base64
import io

# File processing imports
import pandas as pd
from docx import Document as DocxDocument
from pptx import Presentation
import csv
import json

# Disable ChromaDB telemetry using multiple environment variables
os.environ['ANONYMIZED_TELEMETRY'] = 'false'
os.environ['CHROMA_TELEMETRY_DISABLED'] = 'true'
os.environ['CHROMA_TELEMETRY'] = 'false'

# Suppress ChromaDB telemetry error logs
logging.getLogger("chromadb.telemetry.product.posthog").setLevel(logging.CRITICAL)

import chromadb
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_community.document_loaders import PyPDFLoader
from langchain_core.documents import Document
import tempfile
import uuid
from sqlalchemy.orm import Session

from models import RAGDocument

logger = logging.getLogger(__name__)

class RAGService:
    def __init__(self, chroma_db_path: str = None):
        """Initialize RAG service with ChromaDB"""
        if chroma_db_path is None:
            # Point to the chromadb folder in the backend root directory
            chroma_db_path = os.path.join(os.path.dirname(os.path.dirname(os.path.dirname(__file__))), "chromadb")
        
        self.chroma_db_path = chroma_db_path
        self.collection_name = "document_vectors"  # Updated name to reflect multi-file support
        
        # Supported file types and their MIME types
        self.supported_file_types = {
            'pdf': ['application/pdf'],
            'docx': ['application/vnd.openxmlformats-officedocument.wordprocessingml.document'],
            'doc': ['application/msword'],
            'pptx': ['application/vnd.openxmlformats-officedocument.presentationml.presentation'],
            'ppt': ['application/vnd.ms-powerpoint'],
            'xlsx': ['application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'],
            'xls': ['application/vnd.ms-excel'],
            'csv': ['text/csv', 'application/csv'],
            'txt': ['text/plain'],
            'json': ['application/json'],
            'md': ['text/markdown'],
            'odt': ['application/vnd.oasis.opendocument.text'],
            'ods': ['application/vnd.oasis.opendocument.spreadsheet'],
            'odp': ['application/vnd.oasis.opendocument.presentation']
        }
        
        # Initialize ChromaDB client
        self.client = chromadb.PersistentClient(
            path=self.chroma_db_path
        )
        
        # Initialize embeddings model (using a free HuggingFace model)
        self.embeddings = HuggingFaceEmbeddings(
            model_name="sentence-transformers/all-MiniLM-L6-v2",
            model_kwargs={'device': 'cpu'},
            encode_kwargs={'normalize_embeddings': True}
        )
        
        # Initialize text splitter
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1500,
            chunk_overlap=300,
            length_function=len,
            separators=["\n\n", "\n", " ", ""]
        )
        
        # Get or create collection
        try:
            self.collection = self.client.get_collection(name=self.collection_name)
            logger.info(f"✅ Connected to existing ChromaDB collection: {self.collection_name}")
        except Exception:
            # Collection doesn't exist, create it
            self.collection = self.client.create_collection(
                name=self.collection_name
            )
            logger.info(f"✅ Created new ChromaDB collection: {self.collection_name}")
    
    def _generate_file_hash(self, content: bytes) -> str:
        """Generate MD5 hash for file content to check if already processed"""
        return hashlib.md5(content).hexdigest()
    
    def _extract_text_from_pdf_bytes(self, pdf_bytes: bytes, filename: str) -> List[Document]:
        """Extract text from PDF bytes using LangChain's PyPDFLoader"""
        try:
            # Save bytes to temporary file
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as temp_file:
                temp_file.write(pdf_bytes)
                temp_file_path = temp_file.name
            
            try:
                # Load PDF using LangChain's PyPDFLoader
                loader = PyPDFLoader(temp_file_path)
                documents = loader.load()
                
                # Add filename metadata to all documents
                for doc in documents:
                    doc.metadata.update({
                        'filename': filename,
                        'source': filename
                    })
                
                logger.info(f"📄 Extracted {len(documents)} pages from PDF: {filename}")
                return documents
                
            finally:
                # Clean up temporary file
                os.unlink(temp_file_path)
                
        except Exception as e:
            logger.error(f"❌ Failed to extract text from PDF {filename}: {e}")
            return []
    
    def _extract_text_from_docx(self, file_bytes: bytes, filename: str) -> List[Document]:
        """Extract text from DOCX file"""
        try:
            with tempfile.NamedTemporaryFile(delete=False, suffix='.docx') as temp_file:
                temp_file.write(file_bytes)
                temp_file_path = temp_file.name
            
            try:
                doc = DocxDocument(temp_file_path)
                text_content = []
                
                # Extract paragraphs
                for paragraph in doc.paragraphs:
                    if paragraph.text.strip():
                        text_content.append(paragraph.text.strip())
                
                # Extract tables
                for table in doc.tables:
                    for row in table.rows:
                        row_text = []
                        for cell in row.cells:
                            if cell.text.strip():
                                row_text.append(cell.text.strip())
                        if row_text:
                            text_content.append(" | ".join(row_text))
                
                full_text = "\n\n".join(text_content)
                
                document = Document(
                    page_content=full_text,
                    metadata={
                        'filename': filename,
                        'source': filename,
                        'file_type': 'docx'
                    }
                )
                
                logger.info(f"📄 Extracted text from DOCX: {filename}")
                return [document]
                
            finally:
                os.unlink(temp_file_path)
                
        except Exception as e:
            logger.error(f"❌ Failed to extract text from DOCX {filename}: {e}")
            return []
    
    def _extract_text_from_pptx(self, file_bytes: bytes, filename: str) -> List[Document]:
        """Extract text from PPTX file"""
        try:
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as temp_file:
                temp_file.write(file_bytes)
                temp_file_path = temp_file.name
            
            try:
                presentation = Presentation(temp_file_path)
                slides_content = []
                
                for slide_num, slide in enumerate(presentation.slides, 1):
                    slide_text = []
                    
                    # Extract text from shapes
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_text.append(shape.text.strip())
                    
                    if slide_text:
                        slide_content = f"Slide {slide_num}:\n" + "\n".join(slide_text)
                        slides_content.append(slide_content)
                
                full_text = "\n\n".join(slides_content)
                
                document = Document(
                    page_content=full_text,
                    metadata={
                        'filename': filename,
                        'source': filename,
                        'file_type': 'pptx'
                    }
                )
                
                logger.info(f"📄 Extracted text from PPTX: {filename}")
                return [document]
                
            finally:
                os.unlink(temp_file_path)
                
        except Exception as e:
            logger.error(f"❌ Failed to extract text from PPTX {filename}: {e}")
            return []
    
    def _extract_text_from_excel(self, file_bytes: bytes, filename: str) -> List[Document]:
        """Extract text from Excel file (xlsx/xls)"""
        try:
            # Try to read with pandas
            df_dict = pd.read_excel(io.BytesIO(file_bytes), sheet_name=None)
            
            sheets_content = []
            for sheet_name, df in df_dict.items():
                if not df.empty:
                    # Convert DataFrame to string representation
                    sheet_text = f"Sheet: {sheet_name}\n\n"
                    sheet_text += df.to_string(index=False)
                    sheets_content.append(sheet_text)
            
            full_text = "\n\n" + "="*50 + "\n\n".join(sheets_content)
            
            document = Document(
                page_content=full_text,
                metadata={
                    'filename': filename,
                    'source': filename,
                    'file_type': 'excel'
                }
            )
            
            logger.info(f"📄 Extracted text from Excel: {filename}")
            return [document]
            
        except Exception as e:
            logger.error(f"❌ Failed to extract text from Excel {filename}: {e}")
            return []
    
    def _extract_text_from_csv(self, file_bytes: bytes, filename: str) -> List[Document]:
        """Extract text from CSV file"""
        try:
            # Try different encodings
            encodings = ['utf-8', 'latin-1', 'cp1252']
            
            for encoding in encodings:
                try:
                    text_content = file_bytes.decode(encoding)
                    break
                except UnicodeDecodeError:
                    continue
            else:
                raise UnicodeDecodeError("Failed to decode CSV with any supported encoding")
            
            # Parse CSV
            csv_reader = csv.reader(io.StringIO(text_content))
            rows = list(csv_reader)
            
            if rows:
                # Format as readable text
                formatted_text = []
                header = rows[0] if rows else []
                
                formatted_text.append("CSV Data:")
                formatted_text.append("Headers: " + " | ".join(header))
                formatted_text.append("-" * 50)
                
                for row in rows[1:]:
                    if len(row) == len(header):
                        row_text = []
                        for i, cell in enumerate(row):
                            if i < len(header):
                                row_text.append(f"{header[i]}: {cell}")
                        formatted_text.append(" | ".join(row_text))
                
                full_text = "\n".join(formatted_text)
            else:
                full_text = "Empty CSV file"
            
            document = Document(
                page_content=full_text,
                metadata={
                    'filename': filename,
                    'source': filename,
                    'file_type': 'csv'
                }
            )
            
            logger.info(f"📄 Extracted text from CSV: {filename}")
            return [document]
            
        except Exception as e:
            logger.error(f"❌ Failed to extract text from CSV {filename}: {e}")
            return []
    
    def _extract_text_from_json(self, file_bytes: bytes, filename: str) -> List[Document]:
        """Extract text from JSON file"""
        try:
            text_content = file_bytes.decode('utf-8')
            json_data = json.loads(text_content)
            
            # Convert JSON to readable text
            formatted_json = json.dumps(json_data, indent=2, ensure_ascii=False)
            
            document = Document(
                page_content=f"JSON Data from {filename}:\n\n{formatted_json}",
                metadata={
                    'filename': filename,
                    'source': filename,
                    'file_type': 'json'
                }
            )
            
            logger.info(f"📄 Extracted text from JSON: {filename}")
            return [document]
            
        except Exception as e:
            logger.error(f"❌ Failed to extract text from JSON {filename}: {e}")
            return []
    
    def _extract_text_from_text(self, file_bytes: bytes, filename: str) -> List[Document]:
        """Extract text from plain text file"""
        try:
            # Try different encodings
            encodings = ['utf-8', 'latin-1', 'cp1252']
            
            for encoding in encodings:
                try:
                    text_content = file_bytes.decode(encoding)
                    break
                except UnicodeDecodeError:
                    continue
            else:
                raise UnicodeDecodeError("Failed to decode text file with any supported encoding")
            
            document = Document(
                page_content=text_content,
                metadata={
                    'filename': filename,
                    'source': filename,
                    'file_type': 'text'
                }
            )
            
            logger.info(f"📄 Extracted text from text file: {filename}")
            return [document]
            
        except Exception as e:
            logger.error(f"❌ Failed to extract text from text file {filename}: {e}")
            return []
    
    def _get_file_type(self, filename: str, mime_type: str) -> str:
        """Determine file type from filename and MIME type"""
        # Get file extension
        file_ext = filename.lower().split('.')[-1] if '.' in filename else ''
        
        # Map MIME types to file types
        for file_type, mime_types in self.supported_file_types.items():
            if mime_type in mime_types or file_ext == file_type:
                return file_type
        
        return 'unknown'
    
    def _extract_text_from_file(self, file_bytes: bytes, filename: str, mime_type: str) -> List[Document]:
        """Extract text from any supported file type"""
        file_type = self._get_file_type(filename, mime_type)
        
        if file_type == 'pdf':
            return self._extract_text_from_pdf_bytes(file_bytes, filename)
        elif file_type == 'docx':
            return self._extract_text_from_docx(file_bytes, filename)
        elif file_type == 'pptx':
            return self._extract_text_from_pptx(file_bytes, filename)
        elif file_type in ['xlsx', 'xls']:
            return self._extract_text_from_excel(file_bytes, filename)
        elif file_type == 'csv':
            return self._extract_text_from_csv(file_bytes, filename)
        elif file_type == 'json':
            return self._extract_text_from_json(file_bytes, filename)
        elif file_type in ['txt', 'md']:
            return self._extract_text_from_text(file_bytes, filename)
        else:
            logger.warning(f"⚠️ Unsupported file type: {file_type} for file: {filename}")
            return []
    
    def add_file_to_vectorstore(self, file_bytes: bytes, filename: str, mime_type: str, user_email: str, thread_id: int, user_id: int = None, db: Session = None) -> Dict[str, Any]:
        """Add file content to ChromaDB vector store and save metadata to database"""
        try:
            # Generate file hash to check if already processed for this thread
            file_hash = self._generate_file_hash(file_bytes)
            
            # Check if this file has already been processed for this specific thread
            existing_docs = self.collection.get(
                where={
                    "$and": [
                        {"file_hash": file_hash},
                        {"thread_id": thread_id}
                    ]
                }
            )
            
            if existing_docs['ids']:
                file_type = self._get_file_type(filename, mime_type)
                logger.info(f"📋 {file_type.upper()} already exists in vector store for thread {thread_id}: {filename}")
                return {
                    "status": "exists",
                    "message": f"File '{filename}' has already been processed and stored in this conversation.",
                    "file_hash": file_hash,
                    "chunks_count": len(existing_docs['ids'])
                }
            
            # Extract text from file based on type
            documents = self._extract_text_from_file(file_bytes, filename, mime_type)
            
            if not documents:
                return {
                    "status": "error",
                    "message": f"Failed to extract text from file: {filename}"
                }
            
            # Split documents into chunks
            chunks = self.text_splitter.split_documents(documents)
            
            if not chunks:
                return {
                    "status": "error",
                    "message": f"No text chunks generated from file: {filename}"
                }
            
            # Prepare data for ChromaDB
            texts = [chunk.page_content for chunk in chunks]
            metadatas = []
            ids = []
            
            for i, chunk in enumerate(chunks):
                chunk_id = f"{file_hash}_{thread_id}_{i}"  # Include thread_id in chunk_id
                ids.append(chunk_id)
                
                metadata = chunk.metadata.copy()
                metadata.update({
                    'file_hash': file_hash,
                    'chunk_index': i,
                    'user_email': user_email,
                    'thread_id': thread_id,  # Add thread_id to metadata
                    'chunk_id': chunk_id,
                    'mime_type': mime_type
                })
                metadatas.append(metadata)
            
            # Generate embeddings and add to ChromaDB
            embeddings = self.embeddings.embed_documents(texts)
            
            self.collection.add(
                embeddings=embeddings,
                documents=texts,
                metadatas=metadatas,
                ids=ids
            )
            
            # Save metadata to database if db session and user_id are provided
            if db and user_id:
                try:
                    from models import RAGDocument
                    file_type = self._get_file_type(filename, mime_type)
                    
                    rag_doc = RAGDocument(
                        filename=filename,
                        file_hash=file_hash,
                        user_id=user_id,
                        thread_id=thread_id,  # Add thread_id
                        file_size=len(file_bytes),
                        file_type=file_type,
                        chunks_count=len(chunks),
                        total_characters=sum(len(text) for text in texts),
                        status="processed"
                    )
                    
                    db.add(rag_doc)
                    db.commit()
                    logger.info(f"💾 RAG document metadata saved to database: {filename}")
                    
                except Exception as e:
                    logger.error(f"❌ Failed to save RAG document metadata to database: {e}")
                    # Don't fail the whole operation if database save fails
            
            file_type = self._get_file_type(filename, mime_type)
            logger.info(f"✅ Added {len(chunks)} chunks from {file_type.upper()} '{filename}' to vector store for user: {user_email}")
            
            return {
                "status": "success",
                "message": f"Successfully processed and stored file '{filename}' in the knowledge base.",
                "file_hash": file_hash,
                "chunks_count": len(chunks),
                "total_characters": sum(len(text) for text in texts),
                "file_type": file_type
            }
            
        except Exception as e:
            logger.error(f"❌ Failed to add file to vector store: {e}")
            return {
                "status": "error",
                "message": f"Error processing file: {str(e)}"
            }
    
    # Keep backward compatibility for PDF method
    def add_pdf_to_vectorstore(self, pdf_bytes: bytes, filename: str, user_email: str, thread_id: int, user_id: int = None, db: Session = None) -> Dict[str, Any]:
        """Add PDF content to ChromaDB vector store - backward compatibility wrapper"""
        return self.add_file_to_vectorstore(pdf_bytes, filename, 'application/pdf', user_email, thread_id, user_id, db)
    
    def search_relevant_content(self, query: str, user_email: str = None, thread_id: int = None, top_k: int = 5) -> List[Dict[str, Any]]:
        """Search for relevant content in the vector store for specific thread"""
        try:
            # Generate query embedding
            query_embedding = self.embeddings.embed_query(query)
            
            # Build where clause for filtering by user and thread
            where_conditions = []
            if user_email:
                where_conditions.append({"user_email": user_email})
            if thread_id:
                where_conditions.append({"thread_id": thread_id})
            
            where_clause = None
            if where_conditions:
                if len(where_conditions) == 1:
                    where_clause = where_conditions[0]
                else:
                    where_clause = {"$and": where_conditions}
            
            # Search in ChromaDB
            results = self.collection.query(
                query_embeddings=[query_embedding],
                n_results=top_k,
                where=where_clause,
                include=["documents", "metadatas", "distances"]
            )
            
            if not results['documents'][0]:
                logger.info(f"🔍 No relevant content found for query: '{query[:50]}...'")
                return []
            
            # Format results
            relevant_chunks = []
            for i, (doc, metadata, distance) in enumerate(zip(
                results['documents'][0], 
                results['metadatas'][0], 
                results['distances'][0]
            )):
                relevant_chunks.append({
                    'content': doc,
                    'filename': metadata.get('filename', 'Unknown'),
                    'page': metadata.get('page', 'Unknown'),
                    'chunk_index': metadata.get('chunk_index', i),
                    'similarity_score': 1 - distance,  # Convert distance to similarity
                    'metadata': metadata
                })
            
            logger.info(f"🔍 Found {len(relevant_chunks)} relevant chunks for query: '{query[:50]}...'")
            return relevant_chunks
            
        except Exception as e:
            logger.error(f"❌ Failed to search vector store: {e}")
            return []

    def search_in_specific_file(self, query: str, filename: str, user_email: str = None, thread_id: int = None, top_k: int = 5) -> List[Dict[str, Any]]:
        """Search for relevant content in a specific file within a specific thread"""
        try:
            # Generate query embedding
            query_embedding = self.embeddings.embed_query(query)
            
            # Build where clause for filtering by filename, user, and thread
            where_conditions = [{"filename": filename}]
            if user_email:
                where_conditions.append({"user_email": user_email})
            if thread_id:
                where_conditions.append({"thread_id": thread_id})
            
            where_clause = {"$and": where_conditions} if len(where_conditions) > 1 else where_conditions[0]
            
            # Search in ChromaDB
            results = self.collection.query(
                query_embeddings=[query_embedding],
                n_results=top_k,
                where=where_clause,
                include=["documents", "metadatas", "distances"]
            )
            
            if not results['documents'][0]:
                logger.info(f"🔍 No relevant content found in file '{filename}' for query: '{query[:50]}...'")
                return []
            
            # Format results
            relevant_chunks = []
            for i, (doc, metadata, distance) in enumerate(zip(
                results['documents'][0], 
                results['metadatas'][0], 
                results['distances'][0]
            )):
                relevant_chunks.append({
                    'content': doc,
                    'filename': metadata.get('filename', 'Unknown'),
                    'page': metadata.get('page', 'Unknown'),
                    'chunk_index': metadata.get('chunk_index', i),
                    'similarity_score': 1 - distance,  # Convert distance to similarity
                    'metadata': metadata
                })
            
            logger.info(f"🎯 Found {len(relevant_chunks)} relevant chunks in file '{filename}' for query: '{query[:50]}...'")
            return relevant_chunks
            
        except Exception as e:
            logger.error(f"❌ Failed to search in specific file '{filename}': {e}")
            return []
    
    def get_collection_stats(self) -> Dict[str, Any]:
        """Get statistics about the vector store collection"""
        try:
            count = self.collection.count()
            
            # Get all metadata to analyze
            all_docs = self.collection.get(include=["metadatas"])
            
            files = set()
            users = set()
            for metadata in all_docs['metadatas']:
                if 'filename' in metadata:
                    files.add(metadata['filename'])
                if 'user_email' in metadata:
                    users.add(metadata['user_email'])
            
            return {
                "total_chunks": count,
                "unique_files": len(files),
                "unique_users": len(users),
                "files": list(files),
                "collection_name": self.collection_name
            }
            
        except Exception as e:
            logger.error(f"❌ Failed to get collection stats: {e}")
            return {"error": str(e)}
    
    def delete_file_from_vectorstore(self, filename: str, user_email: str = None) -> Dict[str, Any]:
        """Delete all chunks of a specific file from the vector store"""
        try:
            where_clause = {"filename": filename}
            if user_email:
                where_clause["user_email"] = user_email
            
            # Get documents to delete
            docs_to_delete = self.collection.get(where=where_clause)
            
            if not docs_to_delete['ids']:
                return {
                    "status": "not_found",
                    "message": f"No documents found for file: {filename}"
                }
            
            # Delete documents
            self.collection.delete(where=where_clause)
            
            logger.info(f"🗑️ Deleted {len(docs_to_delete['ids'])} chunks for file: {filename}")
            
            return {
                "status": "success",
                "message": f"Successfully deleted {len(docs_to_delete['ids'])} chunks for file: {filename}",
                "deleted_chunks": len(docs_to_delete['ids'])
            }
            
        except Exception as e:
            logger.error(f"❌ Failed to delete file from vector store: {e}")
            return {
                "status": "error",
                "message": f"Error deleting file: {str(e)}"
            }
    
    def reset_collection(self) -> Dict[str, Any]:
        """Reset (clear) the entire collection - USE WITH CAUTION"""
        try:
            # Delete the collection
            self.client.delete_collection(name=self.collection_name)
            
            # Recreate the collection
            self.collection = self.client.create_collection(
                name=self.collection_name
            )
            
            logger.warning(f"⚠️ RESET ChromaDB collection: {self.collection_name}")
            
            return {
                "status": "success",
                "message": f"Successfully reset collection: {self.collection_name}"
            }
            
        except Exception as e:
            logger.error(f"❌ Failed to reset collection: {e}")
            return {
                "status": "error",
                "message": f"Error resetting collection: {str(e)}"
            }
